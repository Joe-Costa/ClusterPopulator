"""File generators for various document formats."""

import asyncio
import csv
import io
import json
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor

import aiofiles
from docx import Document
from docx.shared import Inches, Pt
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle

from .content import ContentGenerator


class FileGenerator:
    """Generate various file types with realistic content."""

    def __init__(self, content_gen: ContentGenerator | None = None):
        """Initialize file generator."""
        self.content = content_gen or ContentGenerator()
        self._executor = ThreadPoolExecutor(max_workers=8)

    async def generate_txt(self, filepath: Path, content_type: str = "memo") -> Path:
        """Generate a plain text file."""
        if content_type == "memo":
            data = self.content.business_memo()
            text = f"""MEMORANDUM

TO: {data['to']}
FROM: {data['from']}
DATE: {data['date']}
SUBJECT: {data['subject']}

{data['body']}
"""
        elif content_type == "notes":
            data = self.content.meeting_notes()
            text = f"""{data['title']}
Date: {data['date']} at {data['time']}

Attendees:
{chr(10).join('- ' + a for a in data['attendees'])}

Agenda:
{chr(10).join('- ' + a for a in data['agenda'])}

Notes:
{data['notes']}

Action Items:
{chr(10).join(f"- {item['task']} (Owner: {item['owner']})" for item in data['action_items'])}
"""
        elif content_type == "log":
            entries = self.content.log_entries(50)
            text = "\n".join(
                f"[{e['timestamp']}] [{e['level']}] [{e['source']}] {e['message']}"
                for e in entries
            )
        else:
            text = self.content.lorem_paragraphs(5)

        async with aiofiles.open(filepath, "w", encoding="utf-8") as f:
            await f.write(text)
        return filepath

    async def generate_json(self, filepath: Path, content_type: str = "config") -> Path:
        """Generate a JSON file."""
        if content_type == "invoice":
            data = self.content.invoice_data()
        elif content_type == "employee":
            data = self.content.employee_record()
        elif content_type == "project":
            data = self.content.project_data()
        elif content_type == "log":
            data = {"entries": self.content.log_entries(30)}
        elif content_type == "config":
            data = {
                "version": "1.0.0",
                "environment": "production",
                "settings": {
                    "debug": False,
                    "log_level": "INFO",
                    "max_connections": 100,
                    "timeout_seconds": 30,
                },
                "features": {
                    "enable_cache": True,
                    "enable_analytics": True,
                    "maintenance_mode": False,
                },
            }
        else:
            data = self.content.project_data()

        async with aiofiles.open(filepath, "w", encoding="utf-8") as f:
            await f.write(json.dumps(data, indent=2))
        return filepath

    async def generate_csv(self, filepath: Path, content_type: str = "data") -> Path:
        """Generate a CSV file."""
        if content_type == "employees":
            headers = ["Employee ID", "Name", "Email", "Department", "Title", "Hire Date", "Salary"]
            rows = []
            for _ in range(50):
                emp = self.content.employee_record()
                rows.append([
                    emp["employee_id"], emp["name"], emp["email"],
                    emp["department"], emp["title"], emp["hire_date"],
                    emp["salary"]
                ])
        elif content_type == "invoices":
            headers = ["Invoice", "Date", "Customer", "Amount", "Status"]
            rows = []
            for _ in range(30):
                inv = self.content.invoice_data()
                rows.append([
                    inv["invoice_number"], inv["date"],
                    inv["customer"]["name"], inv["total"],
                    "Paid" if hash(inv["invoice_number"]) % 2 == 0 else "Pending"
                ])
        else:
            data = self.content.spreadsheet_data(rows=30, cols=6)
            headers = data[0]
            rows = data[1:]

        output = io.StringIO()
        writer = csv.writer(output)
        writer.writerow(headers)
        writer.writerows(rows)

        async with aiofiles.open(filepath, "w", encoding="utf-8", newline="") as f:
            await f.write(output.getvalue())
        return filepath

    def _generate_docx_sync(self, filepath: Path, content_type: str) -> Path:
        """Generate a Word document (sync operation)."""
        doc = Document()

        if content_type == "memo":
            data = self.content.business_memo()
            doc.add_heading("MEMORANDUM", 0)
            doc.add_paragraph(f"TO: {data['to']}")
            doc.add_paragraph(f"FROM: {data['from']}")
            doc.add_paragraph(f"DATE: {data['date']}")
            doc.add_paragraph(f"SUBJECT: {data['subject']}")
            doc.add_paragraph()
            for para in data["body"].split("\n\n"):
                doc.add_paragraph(para)

        elif content_type == "policy":
            data = self.content.policy_document()
            doc.add_heading(data["title"], 0)
            doc.add_paragraph(f"Effective Date: {data['effective_date']}")
            doc.add_paragraph(f"Version: {data['version']}")
            doc.add_paragraph(f"Approved By: {data['approved_by']}")
            doc.add_paragraph()
            doc.add_heading("Purpose", level=1)
            doc.add_paragraph(data["purpose"])
            doc.add_heading("Scope", level=1)
            doc.add_paragraph(data["scope"])
            for section in data["sections"]:
                doc.add_heading(section["heading"], level=1)
                doc.add_paragraph(section["content"])

        elif content_type == "contract":
            data = self.content.contract_data()
            doc.add_heading(data["title"], 0)
            doc.add_paragraph(f"Contract ID: {data['contract_id']}")
            doc.add_paragraph()
            doc.add_paragraph(f"This agreement is entered into between {data['party_a']} (Party A) and {data['party_b']} (Party B).")
            doc.add_paragraph()
            doc.add_paragraph(f"Effective Date: {data['effective_date']}")
            doc.add_paragraph(f"Expiration Date: {data['expiration_date']}")
            doc.add_paragraph(f"Contract Value: ${data['value']:,.2f}")
            doc.add_paragraph()
            doc.add_heading("Terms and Conditions", level=1)
            for para in data["terms"].split("\n\n"):
                doc.add_paragraph(para)

        elif content_type == "report":
            data = self.content.financial_report_data()
            doc.add_heading(f"Financial Report - {data['period']}", 0)
            doc.add_paragraph()
            doc.add_heading("Summary", level=1)
            doc.add_paragraph(f"Revenue: ${data['revenue']:,.2f}")
            doc.add_paragraph(f"Expenses: ${data['expenses']:,.2f}")
            doc.add_paragraph(f"Net Income: ${data['net_income']:,.2f}")
            doc.add_heading("Revenue Breakdown", level=1)
            for cat, val in data["categories"].items():
                doc.add_paragraph(f"{cat}: ${val:,.2f}")
            doc.add_heading("Expense Breakdown", level=1)
            for cat, val in data["expense_breakdown"].items():
                doc.add_paragraph(f"{cat}: ${val:,.2f}")

        else:
            notes = self.content.meeting_notes()
            doc.add_heading(notes["title"], 0)
            doc.add_paragraph(f"Date: {notes['date']} at {notes['time']}")
            doc.add_heading("Attendees", level=1)
            for attendee in notes["attendees"]:
                doc.add_paragraph(attendee, style="List Bullet")
            doc.add_heading("Agenda", level=1)
            for item in notes["agenda"]:
                doc.add_paragraph(item, style="List Bullet")
            doc.add_heading("Notes", level=1)
            doc.add_paragraph(notes["notes"])
            doc.add_heading("Action Items", level=1)
            for item in notes["action_items"]:
                doc.add_paragraph(f"{item['task']} - Owner: {item['owner']}", style="List Bullet")

        doc.save(str(filepath))
        return filepath

    async def generate_docx(self, filepath: Path, content_type: str = "memo") -> Path:
        """Generate a Word document."""
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(
            self._executor, self._generate_docx_sync, filepath, content_type
        )

    def _generate_xlsx_sync(self, filepath: Path, content_type: str) -> Path:
        """Generate an Excel spreadsheet (sync operation)."""
        wb = Workbook()
        ws = wb.active

        header_font = Font(bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")

        if content_type == "financial":
            ws.title = "Financial Report"
            data = self.content.financial_report_data()
            ws["A1"] = f"Financial Report - {data['period']}"
            ws["A1"].font = Font(bold=True, size=14)
            ws.merge_cells("A1:C1")

            ws["A3"] = "Summary"
            ws["A3"].font = Font(bold=True)
            ws["A4"] = "Revenue"
            ws["B4"] = data["revenue"]
            ws["A5"] = "Expenses"
            ws["B5"] = data["expenses"]
            ws["A6"] = "Net Income"
            ws["B6"] = data["net_income"]

            ws["A8"] = "Revenue Breakdown"
            ws["A8"].font = Font(bold=True)
            row = 9
            for cat, val in data["categories"].items():
                ws[f"A{row}"] = cat
                ws[f"B{row}"] = val
                row += 1

        elif content_type == "employees":
            ws.title = "Employee Directory"
            headers = ["Employee ID", "Name", "Email", "Department", "Title", "Hire Date", "Salary"]
            for col, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=header)
                cell.font = header_font
                cell.fill = header_fill

            for row in range(2, 52):
                emp = self.content.employee_record()
                ws.cell(row=row, column=1, value=emp["employee_id"])
                ws.cell(row=row, column=2, value=emp["name"])
                ws.cell(row=row, column=3, value=emp["email"])
                ws.cell(row=row, column=4, value=emp["department"])
                ws.cell(row=row, column=5, value=emp["title"])
                ws.cell(row=row, column=6, value=emp["hire_date"])
                ws.cell(row=row, column=7, value=emp["salary"])

        elif content_type == "invoice":
            ws.title = "Invoice"
            inv = self.content.invoice_data()
            ws["A1"] = "INVOICE"
            ws["A1"].font = Font(bold=True, size=16)
            ws["A3"] = f"Invoice #: {inv['invoice_number']}"
            ws["A4"] = f"Date: {inv['date']}"
            ws["A5"] = f"Due Date: {inv['due_date']}"

            ws["A7"] = "Vendor:"
            ws["A7"].font = Font(bold=True)
            ws["A8"] = inv["vendor"]["name"]
            ws["A9"] = inv["vendor"]["address"].replace("\n", ", ")

            ws["C7"] = "Bill To:"
            ws["C7"].font = Font(bold=True)
            ws["C8"] = inv["customer"]["name"]
            ws["C9"] = inv["customer"]["address"].replace("\n", ", ")

            headers = ["Description", "Quantity", "Unit Price", "Total"]
            for col, header in enumerate(headers, 1):
                cell = ws.cell(row=12, column=col, value=header)
                cell.font = header_font
                cell.fill = header_fill

            row = 13
            for item in inv["items"]:
                ws.cell(row=row, column=1, value=item["description"])
                ws.cell(row=row, column=2, value=item["quantity"])
                ws.cell(row=row, column=3, value=item["unit_price"])
                ws.cell(row=row, column=4, value=item["total"])
                row += 1

            row += 1
            ws.cell(row=row, column=3, value="Subtotal:")
            ws.cell(row=row, column=4, value=inv["subtotal"])
            row += 1
            ws.cell(row=row, column=3, value="Tax:")
            ws.cell(row=row, column=4, value=inv["tax"])
            row += 1
            ws.cell(row=row, column=3, value="Total:")
            ws.cell(row=row, column=3).font = Font(bold=True)
            ws.cell(row=row, column=4, value=inv["total"])
            ws.cell(row=row, column=4).font = Font(bold=True)

        else:
            ws.title = "Data"
            data = self.content.spreadsheet_data(rows=40, cols=8)
            for row_idx, row_data in enumerate(data, 1):
                for col_idx, value in enumerate(row_data, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)
                    if row_idx == 1:
                        cell.font = header_font
                        cell.fill = header_fill

        for column in ws.columns:
            max_length = 0
            column_letter = None
            for cell in column:
                try:
                    if hasattr(cell, 'column_letter'):
                        column_letter = cell.column_letter
                    if cell.value and len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except (TypeError, AttributeError):
                    pass
            if column_letter:
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width

        wb.save(str(filepath))
        return filepath

    async def generate_xlsx(self, filepath: Path, content_type: str = "data") -> Path:
        """Generate an Excel spreadsheet."""
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(
            self._executor, self._generate_xlsx_sync, filepath, content_type
        )

    def _generate_pdf_sync(self, filepath: Path, content_type: str) -> Path:
        """Generate a PDF document (sync operation)."""
        doc = SimpleDocTemplate(str(filepath), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        title_style = ParagraphStyle(
            "CustomTitle",
            parent=styles["Heading1"],
            fontSize=18,
            spaceAfter=20,
        )

        if content_type == "report":
            data = self.content.financial_report_data()
            story.append(Paragraph(f"Financial Report - {data['period']}", title_style))
            story.append(Spacer(1, 12))
            story.append(Paragraph("Summary", styles["Heading2"]))
            story.append(Paragraph(f"Revenue: ${data['revenue']:,.2f}", styles["Normal"]))
            story.append(Paragraph(f"Expenses: ${data['expenses']:,.2f}", styles["Normal"]))
            story.append(Paragraph(f"Net Income: ${data['net_income']:,.2f}", styles["Normal"]))
            story.append(Spacer(1, 12))

            table_data = [["Category", "Amount"]]
            for cat, val in data["categories"].items():
                table_data.append([cat, f"${val:,.2f}"])

            t = Table(table_data)
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4472C4")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("BOTTOMPADDING", (0, 0), (-1, 0), 12),
                ("BACKGROUND", (0, 1), (-1, -1), colors.beige),
                ("GRID", (0, 0), (-1, -1), 1, colors.black),
            ]))
            story.append(t)

        elif content_type == "contract":
            data = self.content.contract_data()
            story.append(Paragraph(data["title"], title_style))
            story.append(Paragraph(f"Contract ID: {data['contract_id']}", styles["Normal"]))
            story.append(Spacer(1, 12))
            story.append(Paragraph(
                f"This agreement is entered into between <b>{data['party_a']}</b> (Party A) "
                f"and <b>{data['party_b']}</b> (Party B).",
                styles["Normal"]
            ))
            story.append(Spacer(1, 12))
            story.append(Paragraph(f"Effective Date: {data['effective_date']}", styles["Normal"]))
            story.append(Paragraph(f"Expiration Date: {data['expiration_date']}", styles["Normal"]))
            story.append(Paragraph(f"Contract Value: ${data['value']:,.2f}", styles["Normal"]))
            story.append(Spacer(1, 12))
            story.append(Paragraph("Terms and Conditions", styles["Heading2"]))
            for para in data["terms"].split("\n\n"):
                story.append(Paragraph(para, styles["Normal"]))
                story.append(Spacer(1, 6))

        elif content_type == "policy":
            data = self.content.policy_document()
            story.append(Paragraph(data["title"], title_style))
            story.append(Paragraph(f"Effective Date: {data['effective_date']}", styles["Normal"]))
            story.append(Paragraph(f"Version: {data['version']}", styles["Normal"]))
            story.append(Spacer(1, 12))
            story.append(Paragraph("Purpose", styles["Heading2"]))
            story.append(Paragraph(data["purpose"], styles["Normal"]))
            story.append(Spacer(1, 6))
            story.append(Paragraph("Scope", styles["Heading2"]))
            story.append(Paragraph(data["scope"], styles["Normal"]))
            story.append(Spacer(1, 12))
            for section in data["sections"]:
                story.append(Paragraph(section["heading"], styles["Heading2"]))
                story.append(Paragraph(section["content"], styles["Normal"]))
                story.append(Spacer(1, 6))

        elif content_type == "invoice":
            inv = self.content.invoice_data()
            story.append(Paragraph("INVOICE", title_style))
            story.append(Paragraph(f"Invoice #: {inv['invoice_number']}", styles["Normal"]))
            story.append(Paragraph(f"Date: {inv['date']}", styles["Normal"]))
            story.append(Paragraph(f"Due Date: {inv['due_date']}", styles["Normal"]))
            story.append(Spacer(1, 12))

            story.append(Paragraph(f"<b>Vendor:</b> {inv['vendor']['name']}", styles["Normal"]))
            story.append(Paragraph(f"<b>Customer:</b> {inv['customer']['name']}", styles["Normal"]))
            story.append(Spacer(1, 12))

            table_data = [["Description", "Qty", "Unit Price", "Total"]]
            for item in inv["items"]:
                table_data.append([
                    item["description"][:30],
                    str(item["quantity"]),
                    f"${item['unit_price']:.2f}",
                    f"${item['total']:.2f}"
                ])
            table_data.append(["", "", "Subtotal:", f"${inv['subtotal']:.2f}"])
            table_data.append(["", "", "Tax:", f"${inv['tax']:.2f}"])
            table_data.append(["", "", "Total:", f"${inv['total']:.2f}"])

            t = Table(table_data, colWidths=[200, 50, 80, 80])
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4472C4")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                ("ALIGN", (1, 0), (-1, -1), "RIGHT"),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("BOTTOMPADDING", (0, 0), (-1, 0), 12),
                ("GRID", (0, 0), (-1, -4), 1, colors.black),
                ("FONTNAME", (2, -1), (3, -1), "Helvetica-Bold"),
            ]))
            story.append(t)

        else:
            memo = self.content.business_memo()
            story.append(Paragraph("MEMORANDUM", title_style))
            story.append(Paragraph(f"<b>TO:</b> {memo['to']}", styles["Normal"]))
            story.append(Paragraph(f"<b>FROM:</b> {memo['from']}", styles["Normal"]))
            story.append(Paragraph(f"<b>DATE:</b> {memo['date']}", styles["Normal"]))
            story.append(Paragraph(f"<b>SUBJECT:</b> {memo['subject']}", styles["Normal"]))
            story.append(Spacer(1, 20))
            for para in memo["body"].split("\n\n"):
                story.append(Paragraph(para, styles["Normal"]))
                story.append(Spacer(1, 12))

        doc.build(story)
        return filepath

    async def generate_pdf(self, filepath: Path, content_type: str = "memo") -> Path:
        """Generate a PDF document."""
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(
            self._executor, self._generate_pdf_sync, filepath, content_type
        )

    def _generate_pptx_sync(self, filepath: Path, content_type: str) -> Path:
        """Generate a PowerPoint presentation (sync operation)."""
        prs = Presentation()
        slides_content = self.content.presentation_content(slides=10)

        for slide_data in slides_content:
            if slide_data["type"] == "title":
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                title_box = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(2.5), PptxInches(9), PptxInches(1.5)
                )
                tf = title_box.text_frame
                tf.text = slide_data["title"]
                tf.paragraphs[0].font.size = PptxPt(44)
                tf.paragraphs[0].font.bold = True

                if "subtitle" in slide_data:
                    subtitle_box = slide.shapes.add_textbox(
                        PptxInches(0.5), PptxInches(4), PptxInches(9), PptxInches(1)
                    )
                    stf = subtitle_box.text_frame
                    stf.text = slide_data["subtitle"]
                    stf.paragraphs[0].font.size = PptxPt(24)

            elif slide_data["type"] == "bullet":
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                title_box = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(0.5), PptxInches(9), PptxInches(1)
                )
                tf = title_box.text_frame
                tf.text = slide_data["title"]
                tf.paragraphs[0].font.size = PptxPt(32)
                tf.paragraphs[0].font.bold = True

                content_box = slide.shapes.add_textbox(
                    PptxInches(0.75), PptxInches(1.75), PptxInches(8.5), PptxInches(5)
                )
                ctf = content_box.text_frame
                for i, bullet in enumerate(slide_data["bullets"]):
                    if i == 0:
                        ctf.text = bullet
                    else:
                        p = ctf.add_paragraph()
                        p.text = bullet
                    ctf.paragraphs[i].font.size = PptxPt(18)
                    ctf.paragraphs[i].space_before = PptxPt(12)

            elif slide_data["type"] == "content":
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                title_box = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(0.5), PptxInches(9), PptxInches(1)
                )
                tf = title_box.text_frame
                tf.text = slide_data["title"]
                tf.paragraphs[0].font.size = PptxPt(32)
                tf.paragraphs[0].font.bold = True

                content_box = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(1.75), PptxInches(9), PptxInches(5)
                )
                ctf = content_box.text_frame
                ctf.word_wrap = True
                ctf.text = slide_data["content"]
                ctf.paragraphs[0].font.size = PptxPt(16)

            elif slide_data["type"] == "two_column":
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                title_box = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(0.5), PptxInches(9), PptxInches(1)
                )
                tf = title_box.text_frame
                tf.text = slide_data["title"]
                tf.paragraphs[0].font.size = PptxPt(32)
                tf.paragraphs[0].font.bold = True

                left_box = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(1.75), PptxInches(4.25), PptxInches(5)
                )
                ltf = left_box.text_frame
                for i, item in enumerate(slide_data["left"]):
                    if i == 0:
                        ltf.text = item
                    else:
                        p = ltf.add_paragraph()
                        p.text = item
                    ltf.paragraphs[i].font.size = PptxPt(16)

                right_box = slide.shapes.add_textbox(
                    PptxInches(5.25), PptxInches(1.75), PptxInches(4.25), PptxInches(5)
                )
                rtf = right_box.text_frame
                for i, item in enumerate(slide_data["right"]):
                    if i == 0:
                        rtf.text = item
                    else:
                        p = rtf.add_paragraph()
                        p.text = item
                    rtf.paragraphs[i].font.size = PptxPt(16)

        prs.save(str(filepath))
        return filepath

    async def generate_pptx(self, filepath: Path, content_type: str = "presentation") -> Path:
        """Generate a PowerPoint presentation."""
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(
            self._executor, self._generate_pptx_sync, filepath, content_type
        )

    async def generate_xml(self, filepath: Path, content_type: str = "config") -> Path:
        """Generate an XML file."""
        if content_type == "config":
            xml_content = """<?xml version="1.0" encoding="UTF-8"?>
<configuration>
    <application>
        <name>Enterprise Application</name>
        <version>2.1.0</version>
        <environment>production</environment>
    </application>
    <database>
        <host>db.internal.company.com</host>
        <port>5432</port>
        <name>enterprise_db</name>
        <pool_size>20</pool_size>
    </database>
    <logging>
        <level>INFO</level>
        <file>/var/log/app/application.log</file>
        <max_size_mb>100</max_size_mb>
        <backup_count>5</backup_count>
    </logging>
    <features>
        <cache enabled="true" ttl="3600"/>
        <analytics enabled="true"/>
        <maintenance_mode enabled="false"/>
    </features>
</configuration>"""
        elif content_type == "data":
            employees = [self.content.employee_record() for _ in range(10)]
            emp_xml = "\n".join(f"""    <employee id="{e['employee_id']}">
        <name>{e['name']}</name>
        <email>{e['email']}</email>
        <department>{e['department']}</department>
        <title>{e['title']}</title>
    </employee>""" for e in employees)
            xml_content = f"""<?xml version="1.0" encoding="UTF-8"?>
<employees>
{emp_xml}
</employees>"""
        else:
            project = self.content.project_data()
            xml_content = f"""<?xml version="1.0" encoding="UTF-8"?>
<project id="{project['project_id']}">
    <name>{project['name']}</name>
    <description>{project['description']}</description>
    <status>{project['status']}</status>
    <start_date>{project['start_date']}</start_date>
    <target_date>{project['target_date']}</target_date>
    <budget>{project['budget']}</budget>
    <manager>{project['manager']}</manager>
    <team>
        {''.join(f"<member>{m}</member>" for m in project['team_members'])}
    </team>
</project>"""

        async with aiofiles.open(filepath, "w", encoding="utf-8") as f:
            await f.write(xml_content)
        return filepath

    async def generate_html(self, filepath: Path, content_type: str = "report") -> Path:
        """Generate an HTML file."""
        if content_type == "report":
            data = self.content.financial_report_data()
            html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Financial Report - {data['period']}</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 40px; }}
        h1 {{ color: #333; }}
        table {{ border-collapse: collapse; width: 100%; margin: 20px 0; }}
        th, td {{ border: 1px solid #ddd; padding: 12px; text-align: left; }}
        th {{ background-color: #4472C4; color: white; }}
        tr:nth-child(even) {{ background-color: #f2f2f2; }}
        .summary {{ background-color: #f9f9f9; padding: 20px; border-radius: 8px; margin: 20px 0; }}
    </style>
</head>
<body>
    <h1>Financial Report - {data['period']}</h1>
    <div class="summary">
        <h2>Summary</h2>
        <p><strong>Revenue:</strong> ${data['revenue']:,.2f}</p>
        <p><strong>Expenses:</strong> ${data['expenses']:,.2f}</p>
        <p><strong>Net Income:</strong> ${data['net_income']:,.2f}</p>
    </div>
    <h2>Revenue Breakdown</h2>
    <table>
        <tr><th>Category</th><th>Amount</th></tr>
        {''.join(f"<tr><td>{cat}</td><td>${val:,.2f}</td></tr>" for cat, val in data['categories'].items())}
    </table>
    <h2>Expense Breakdown</h2>
    <table>
        <tr><th>Category</th><th>Amount</th></tr>
        {''.join(f"<tr><td>{cat}</td><td>${val:,.2f}</td></tr>" for cat, val in data['expense_breakdown'].items())}
    </table>
</body>
</html>"""
        else:
            policy = self.content.policy_document()
            sections_html = "\n".join(
                f"<h2>{s['heading']}</h2><p>{s['content']}</p>"
                for s in policy['sections']
            )
            html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{policy['title']}</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }}
        h1 {{ color: #333; border-bottom: 2px solid #4472C4; padding-bottom: 10px; }}
        h2 {{ color: #4472C4; margin-top: 30px; }}
        .meta {{ background-color: #f5f5f5; padding: 15px; border-radius: 5px; margin: 20px 0; }}
    </style>
</head>
<body>
    <h1>{policy['title']}</h1>
    <div class="meta">
        <p><strong>Effective Date:</strong> {policy['effective_date']}</p>
        <p><strong>Version:</strong> {policy['version']}</p>
        <p><strong>Approved By:</strong> {policy['approved_by']}</p>
    </div>
    <h2>Purpose</h2>
    <p>{policy['purpose']}</p>
    <h2>Scope</h2>
    <p>{policy['scope']}</p>
    {sections_html}
</body>
</html>"""

        async with aiofiles.open(filepath, "w", encoding="utf-8") as f:
            await f.write(html_content)
        return filepath

    async def generate_md(self, filepath: Path, content_type: str = "notes") -> Path:
        """Generate a Markdown file."""
        if content_type == "notes":
            notes = self.content.meeting_notes()
            md_content = f"""# {notes['title']}

**Date:** {notes['date']} at {notes['time']}

## Attendees

{chr(10).join('- ' + a for a in notes['attendees'])}

## Agenda

{chr(10).join('- ' + a for a in notes['agenda'])}

## Notes

{notes['notes']}

## Action Items

{chr(10).join(f"- [ ] {item['task']} - **Owner:** {item['owner']}" for item in notes['action_items'])}
"""
        elif content_type == "project":
            project = self.content.project_data()
            md_content = f"""# {project['name']}

**Project ID:** {project['project_id']}

## Overview

{project['description']}

## Details

| Field | Value |
|-------|-------|
| Status | {project['status']} |
| Start Date | {project['start_date']} |
| Target Date | {project['target_date']} |
| Budget | ${project['budget']:,.2f} |
| Manager | {project['manager']} |

## Team Members

{chr(10).join('- ' + m for m in project['team_members'])}
"""
        else:
            policy = self.content.policy_document()
            sections_md = "\n\n".join(
                f"## {s['heading']}\n\n{s['content']}"
                for s in policy['sections']
            )
            md_content = f"""# {policy['title']}

**Effective Date:** {policy['effective_date']}
**Version:** {policy['version']}
**Approved By:** {policy['approved_by']}

## Purpose

{policy['purpose']}

## Scope

{policy['scope']}

{sections_md}
"""

        async with aiofiles.open(filepath, "w", encoding="utf-8") as f:
            await f.write(md_content)
        return filepath

    def shutdown(self):
        """Shutdown the executor."""
        self._executor.shutdown(wait=True)
